"""Parse uploaded files into text/image artifacts for multimodal chat."""

from __future__ import annotations

import csv
import io
from pathlib import Path

from fastapi import UploadFile
from PyPDF2 import PdfReader

from ai_backend_v2.models.schemas import UploadedFileStatus
from ai_backend_v2.services.upload_session_service import UploadedArtifact
from ai_backend_v2.utils.logger import get_logger

logger = get_logger(__name__)

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
TEXT_EXTS = {".txt", ".md", ".json", ".xml", ".yaml", ".yml"}


def _safe_decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _parse_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip()


def _parse_csv(data: bytes) -> str:
    text = _safe_decode(data)
    rows = list(csv.reader(io.StringIO(text)))
    lines = [", ".join(cell.strip() for cell in row) for row in rows]
    return "\n".join(lines).strip()


def _parse_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for .docx parsing") from exc
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def _parse_xlsx(data: bytes) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for .xlsx parsing") from exc
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines).strip()


def _parse_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx parsing") from exc
    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {idx}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                lines.append(text.strip())
    return "\n".join(line for line in lines if line).strip()


def parse_file_bytes_to_text(name: str, mime: str, data: bytes) -> str:
    ext = Path(name).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(data)
    if ext in TEXT_EXTS:
        return _safe_decode(data).strip()
    if ext == ".csv":
        return _parse_csv(data)
    if ext == ".docx":
        return _parse_docx(data)
    if ext in {".xlsx", ".xlsm"}:
        return _parse_xlsx(data)
    if ext == ".pptx":
        return _parse_pptx(data)
    raise ValueError(f"Unsupported file type: {ext or mime}")


async def parse_uploads(files: list[UploadFile]) -> tuple[list[UploadedArtifact], list[UploadedFileStatus]]:
    artifacts: list[UploadedArtifact] = []
    statuses: list[UploadedFileStatus] = []

    for file in files:
        name = file.filename or "upload"
        mime = file.content_type or "application/octet-stream"
        ext = Path(name).suffix.lower()
        try:
            data = await file.read()
            if not data:
                statuses.append(UploadedFileStatus(name=name, mime_type=mime, status="failed", detail="Empty file"))
                continue

            if ext in IMAGE_EXTS or mime.startswith("image/"):
                artifacts.append(UploadedArtifact(name=name, mime_type=mime, kind="image", image_bytes=data))
                statuses.append(UploadedFileStatus(name=name, mime_type=mime, status="accepted"))
                continue

            text = parse_file_bytes_to_text(name, mime, data)
            if not text:
                statuses.append(
                    UploadedFileStatus(
                        name=name,
                        mime_type=mime,
                        status="failed",
                        detail="No readable text extracted",
                    )
                )
                continue

            artifacts.append(UploadedArtifact(name=name, mime_type=mime, kind="text", text=text))
            statuses.append(UploadedFileStatus(name=name, mime_type=mime, status="accepted"))
        except ValueError as exc:
            statuses.append(
                UploadedFileStatus(
                    name=name,
                    mime_type=mime,
                    status="unsupported",
                    detail=str(exc),
                )
            )
        except Exception as exc:
            logger.exception("Failed to parse upload: %s", name)
            statuses.append(UploadedFileStatus(name=name, mime_type=mime, status="failed", detail=str(exc)))

    return artifacts, statuses
